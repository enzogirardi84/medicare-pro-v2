"""
Genera presentacion institucional PPTX de MediCare Enterprise PRO.
Con logo, gradientes suaves, tipografia limpia, 10 slides.
"""

import os, tempfile
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

REPO = Path(__file__).resolve().parent.parent
LOGO = REPO / "assets" / "logo_medicare_pro.jpeg"
PPTX = REPO / "marketing" / "Presentacion_MediCare_Enterprise_PRO.pptx"
PDF  = REPO / "marketing" / "Presentacion_MediCare_Enterprise_PRO.pdf"

# Paleta
C   = {"teal": RGBColor(20,184,166), "teal_l":RGBColor(45,212,191),
       "blue": RGBColor(59,130,246), "blue_l":RGBColor(96,165,250),
       "violet":RGBColor(139,92,246), "gold":RGBColor(251,191,36),
       "red":RGBColor(248,113,113),
       "d1":RGBColor(5,8,18), "d2":RGBColor(12,18,34), "d3":RGBColor(18,26,42),
       "card":RGBColor(14,20,34), "card_b":RGBColor(50,60,80),
       "w":RGBColor(255,255,255), "g":RGBColor(148,163,184), "g2":RGBColor(200,210,225),
       "line":RGBColor(55,65,81)}

W, H = Inches(13.333), Inches(7.5)
prs = Presentation(); prs.slide_width = W; prs.slide_height = H

def R(s, l, t, w, h, c):   # rect lleno
    x = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    x.fill.solid(); x.fill.fore_color.rgb = c; x.line.fill.background(); return x
def RR(s, l, t, w, h, c, b=None):  # rounded rect
    x = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    x.fill.solid(); x.fill.fore_color.rgb = c; x.line.fill.background()
    if b: x.line.color.rgb = b; x.line.width = Pt(0.5)
    return x
def O(s, l, t, w, h, c):  # oval
    x = s.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    x.fill.solid(); x.fill.fore_color.rgb = c; x.line.fill.background(); return x
def TX(s, l, t, w, h): return s.shapes.add_textbox(l, t, w, h)

def T(obj, txt, sz=18, b=False, c=None, al=PP_ALIGN.LEFT, fn="Calibri"):
    tf = obj if hasattr(obj,'paragraphs') else obj.text_frame
    p = tf.paragraphs[0] if len(tf.paragraphs)==1 and tf.paragraphs[0].text=="" else tf.add_paragraph()
    p.text = txt; p.font.size=Pt(sz); p.font.bold=b; p.font.name=fn
    p.font.color.rgb = c or C["w"]; p.alignment = al; p.space_after=Pt(6); return p

def B(tf, txt, sz=16, c=None):
    return T(tf, f"  \u25b8  {txt}", sz, False, c or C["g"])

def logo_c(slide, l, t, sz):
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), l, t, sz, sz)
    O(slide, l-Inches(0.04), t-Inches(0.04), sz+Inches(0.08), sz+Inches(0.08), C["teal"])

# ─── 1: PORTADA ──────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6]); R(s,0,0,W,H,C["d1"]); R(s,0,0,W,Inches(0.06),C["teal"])
O(s,Inches(-2),Inches(-1.5),Inches(6),Inches(6),C["blue"])
O(s,Inches(10),Inches(4.5),Inches(4),Inches(4),C["teal"])
for x in range(6):
    R(s, Inches(0.8+x*2), Inches(1.0+x*0.12), Inches(0.015), Inches(6), C["blue"])
logo_c(s, Inches(0.8), Inches(1.0), Inches(1.3))
tx = TX(s, Inches(2.5), Inches(1.0), Inches(9.5), Inches(1.2))
T(tx, "MediCare Enterprise PRO", 50, True, C["w"])
T(tx, "Plataforma integral de gestión sanitaria", 24, False, C["teal_l"])
R(s, Inches(2.5), Inches(2.9), Inches(3.5), Inches(0.04), C["teal"])
tx2 = TX(s, Inches(2.5), Inches(3.2), Inches(9.5), Inches(2.5))
for item in ["Historia clínica · Recetas con firma · Visitas GPS",
             "Emergencias · Telemedicina · App paciente · Chatbot IA",
             "RRHH · Caja · Inventario · Facturación · Auditoría legal"]:
    B(tx2, item, 17, C["g2"])
# Contacto
RR(s, Inches(2.5), Inches(5.8), Inches(8), Inches(1.2), C["card"], C["teal"])
t3 = TX(s, Inches(3.0), Inches(5.9), Inches(7.2), Inches(1.0))
T(t3, "Enzo Girardi  ·  +54 9 358 430 2024  ·  enzogirardi84@gmail.com", 14, False, C["g"])
T(t3, "Dario Lanfranco  ·  +54 9 358 420 1263  ·  dariolanfrancoruffener@gmail.com", 14, False, C["g"])
R(s,0,Inches(7.44),W,Inches(0.06),C["teal"])

# ─── 2: PROBLEMA ──────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6]); R(s,0,0,W,H,C["d1"]); R(s,0,0,W,Inches(0.06),C["red"])
tx = TX(s, Inches(0.8), Inches(0.4), Inches(11), Inches(1))
T(tx, "¿Cómo opera hoy su institución?", 38, True, C["w"])
T(tx, "Los silos de información cuestan tiempo, dinero y seguridad legal", 18, False, C["g"])
R(s, Inches(0.8), Inches(1.5), Inches(3), Inches(0.04), C["gold"])
for i,(tit,desc) in enumerate([
    ("Sin trazabilidad", "Papel, PDFs sueltos, planillas paralelas."),
    ("Sin métricas", "Agenda en Excel, sin KPIs ni alertas."),
    ("Sin control", "Visitas sin GPS, horarios no verificables."),
    ("Sin validez legal", "Recetas a mano, documentación informal."),
    ("Sin respaldo", "WhatsApp como registro oficial."),
    ("Doble carga", "Facturación en sistema aparte, errores."),
]):
    r,c_ = i//3,i%3; x = Inches(0.5 + c_*4.2); y = Inches(1.8+r*2.6)
    RR(s, x, y, Inches(3.9), Inches(2.2), C["card"], C["line"])
    # Icono circulo
    O(s, x+Inches(0.3), y+Inches(0.2), Inches(0.8), Inches(0.8), C["red"])
    t2 = TX(s, x+Inches(1.3), y+Inches(0.3), Inches(2.3), Inches(0.5))
    T(t2, tit, 18, True, C["w"])
    t3 = TX(s, x+Inches(0.3), y+Inches(1.2), Inches(3.3), Inches(0.8))
    T(t3, desc, 14, False, C["g"])
R(s,0,Inches(7.44),W,Inches(0.06),C["red"])

# ─── 3: SOLUCIÓN (módulos) ────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6]); R(s,0,0,W,H,C["d1"]); R(s,0,0,W,Inches(0.06),C["teal"])
tx = TX(s, Inches(0.6), Inches(0.3), Inches(12), Inches(1))
T(tx, "Una plataforma, todas las áreas", 38, True, C["w"])
T(tx, "35+ módulos integrados en un mismo entorno web seguro", 18, False, C["teal_l"])
R(s, Inches(0.6), Inches(1.3), Inches(3), Inches(0.04), C["teal"])
for i,(nom,desc) in enumerate([
    ("Dashboard", "KPIs, heatmap, mapa GPS"),
    ("Agenda", "Visitas, profesionales, turnos"),
    ("Historia Clínica", "Vitales, evolución, escalas"),
    ("Recetas", "Firma digital, vademécum"),
    ("Emergencias", "Triage, alertas, traslado"),
    ("Telemedicina", "Sala remota, app paciente"),
    ("Visitas GPS", "Fichada geolocalizada"),
    ("Farmacopea", "Interacciones, 50+ fármacos"),
    ("RRHH", "Fichajes, presentismo"),
    ("Caja", "Cobros, comprobantes"),
    ("Facturación", "Exportación contable"),
    ("Auditoría", "Trazabilidad legal total"),
    ("Inventario", "Stock, insumos, pedidos"),
    ("Chatbot IA", "Asistente clínico"),
]):
    r,c_ = i//4,i%4; x = Inches(0.4+c_*3.2); y = Inches(1.6+r*1.4)
    RR(s, x, y, Inches(3.0), Inches(1.2), C["card"], C["line"])
    t = TX(s, x+Inches(0.2), y+Inches(0.15), Inches(2.6), Inches(0.9))
    T(t, nom, 16, True, C["teal_l"]); T(t, desc, 11, False, C["g"])
tx = TX(s, Inches(0.6), Inches(7.0), Inches(12), Inches(0.4))
T(tx, "Implementacion en 24-72 hs  ·  Sin instalar software  ·  Acceso por roles  ·  Cifrado HTTPS", 14, False, C["g"])
R(s,0,Inches(7.44),W,Inches(0.06),C["teal"])

# ─── 4: DASHBOARD ──────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6]); R(s,0,0,W,H,C["d1"]); R(s,0,0,W,Inches(0.06),C["blue"])
tx = TX(s, Inches(0.8), Inches(0.4), Inches(11), Inches(1))
T(tx, "Dashboard Ejecutivo", 38, True, C["w"])
T(tx, "Toda su operación en tiempo real, en una sola pantalla", 18, False, C["blue_l"])
R(s, Inches(0.8), Inches(1.5), Inches(3), Inches(0.04), C["blue"])
for i,(tit,desc) in enumerate([
    ("KPIs en tiempo real", "Pacientes activos, visitas del dia, urgencias, ingresos del mes."),
    ("Mapa geográfico de visitas", "Ubicacion de cada profesional con fichada GPS y control horario."),
    ("Calendario heatmap", "30 dias de actividad con picos de demanda y capacidad ociosa."),
    ("Graficos evolutivos", "Tendencia semanal de metricas clave: visitas, facturacion, stock.")]):
    y = Inches(1.8+i*1.3)
    RR(s, Inches(0.8), y, Inches(8), Inches(1.1), C["card"], C["line"])
    t = TX(s, Inches(1.2), y+Inches(0.1), Inches(7.3), Inches(0.9))
    T(t, tit, 18, True, C["w"]); T(t, desc, 13, False, C["g"])
# Metricas
for i,(v,l) in enumerate([("35+","Modulos"),("24-72h","Implementacion"),("Web+Movil","Acceso")]):
    y = Inches(1.8+i*1.3)
    RR(s, Inches(9.2), y, Inches(3.6), Inches(1.1), C["card"], C["teal"])
    t = TX(s, Inches(9.5), y+Inches(0.1), Inches(3.0), Inches(0.9))
    T(t, v, 28, True, C["teal_l"], PP_ALIGN.CENTER); T(t, l, 13, False, C["g"], PP_ALIGN.CENTER)
R(s,0,Inches(7.44),W,Inches(0.06),C["blue"])

# ─── 5: VISITAS GPS ────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6]); R(s,0,0,W,H,C["d1"]); R(s,0,0,W,Inches(0.06),C["gold"])
tx = TX(s, Inches(0.8), Inches(0.4), Inches(11), Inches(1))
T(tx, "Visitas con Fichada GPS", 38, True, C["w"])
T(tx, "Control geográfico, horario y documental de cada intervención", 18, False, C["gold"])
R(s, Inches(0.8), Inches(1.5), Inches(3), Inches(0.04), C["gold"])
for i,it in enumerate([
    "Geolocalización de llegada y salida en cada visita",
    "Control horario con alertas de atraso y ausencia",
    "Plan de visitas asignado por profesional y zona",
    "Documentación asociada: evolución, fotos, formularios",
    "Reporte mensual de cumplimiento con métricas"]):
    y = Inches(1.8+i*0.95)
    RR(s, Inches(0.8), y, Inches(7.5), Inches(0.8), C["card"], C["line"])
    t = TX(s, Inches(1.2), y+Inches(0.1), Inches(6.8), Inches(0.6))
    T(t, f"  {i+1}.  {it}", 16, False, C["g2"])
# Panel
for i,it in enumerate(["Sin planillas paralelas","Desvios cero no reportados",
                        "Respaldo ante auditoria","Metrica real de productividad"]):
    y = Inches(1.8+i*1.3)
    RR(s, Inches(9.0), y, Inches(3.8), Inches(1.1), C["card"], C["teal"])
    t = TX(s, Inches(9.3), y+Inches(0.25), Inches(3.2), Inches(0.6))
    T(t, f"  ✓  {it}", 15, False, C["teal_l"])
R(s,0,Inches(7.44),W,Inches(0.06),C["gold"])

# ─── 6: SEGURIDAD ──────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6]); R(s,0,0,W,H,C["d1"]); R(s,0,0,W,Inches(0.06),C["teal"])
tx = TX(s, Inches(0.8), Inches(0.4), Inches(11), Inches(1))
T(tx, "Seguridad y Cumplimiento", 38, True, C["w"])
T(tx, "Sus datos protegidos en todo momento", 18, False, C["teal_l"])
R(s, Inches(0.8), Inches(1.5), Inches(3), Inches(0.04), C["teal"])
for i,(t,desc) in enumerate([
    ("Cifrado HTTPS", "Toda comunicación viaja cifrada extremo a extremo."),
    ("Acceso por roles", "Cada usuario accede solo a su nivel de responsabilidad."),
    ("2FA opcional", "Doble factor de autenticación por correo."),
    ("Auditoría legal", "Trazabilidad completa de cada acción en el sistema."),
    ("Backup automatizado", "Respaldo diario con redundancia geográfica."),
    ("Protección XSS", "Sanitización de datos y rate limiting contra fuerza bruta.")]):
    r,c_ = i//2,i%2; x = Inches(0.6+c_*6.2); y = Inches(1.8+r*1.7)
    RR(s, x, y, Inches(5.8), Inches(1.4), C["card"], C["line"])
    t2 = TX(s, x+Inches(0.3), y+Inches(0.2), Inches(5.2), Inches(1.0))
    T(t2, t, 18, True, C["w"]); T(t2, desc, 13, False, C["g"])
R(s,0,Inches(7.44),W,Inches(0.06),C["teal"])

# ─── 7: BENEFICIOS ─────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6]); R(s,0,0,W,H,C["d1"]); R(s,0,0,W,Inches(0.06),C["violet"])
tx = TX(s, Inches(0.8), Inches(0.4), Inches(11), Inches(1))
T(tx, "Beneficios", 38, True, C["w"])
R(s, Inches(0.8), Inches(1.3), Inches(3), Inches(0.04), C["violet"])
for i,(ic,tit,desc) in enumerate([
    ("rapido", "Implementación rápida", "Operativa en 24-72 horas. Sin infraestructura propia."),
    ("mobile", "Acceso móvil total", "Funciona en celular y escritorio. Sin instalar nada."),
    ("integrado", "Todo en uno", "35+ módulos integrados. Sin sistemas paralelos."),
    ("legal", "Respaldo legal", "PDF, trazabilidad, auditoría lista para presentar."),
    ("ahorro", "Menor costo", "Elimina planillas, doble carga, errores administrativos."),
    ("ia", "Inteligencia artificial", "Chatbot clínico, alertas, calculadora de dosis."),
    ("capacita", "Capacitación incluida", "Videollamada guiada + manual interactivo."),
    ("soporte", "Soporte directo", "WhatsApp y email con el equipo de desarrollo.")]):
    r,c_ = i//2,i%2; x = Inches(0.6+c_*6.2); y = Inches(1.6+r*1.35)
    RR(s, x, y, Inches(5.8), Inches(1.2), C["card"], C["line"])
    t2 = TX(s, x+Inches(0.3), y+Inches(0.12), Inches(5.2), Inches(0.9))
    T(t2, tit, 18, True, C["violet"]); T(t2, desc, 13, False, C["g"])
R(s,0,Inches(7.44),W,Inches(0.06),C["violet"])

# ─── 8: PLANES ─────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6]); R(s,0,0,W,H,C["d1"]); R(s,0,0,W,Inches(0.06),C["gold"])
tx = TX(s, Inches(0.8), Inches(0.4), Inches(11), Inches(1))
T(tx, "Planes", 38, True, C["w"]); T(tx, "Elegí el plan ideal para tu institución", 18, False, C["gold"])
R(s, Inches(0.8), Inches(1.5), Inches(3), Inches(0.04), C["gold"])
for i,(nom,lim,color,items) in enumerate([
    ("Básico", "Hasta 10 profesionales", C["teal"],
     ["Dashboard ejecutivo","Agenda de visitas","Historia clínica","Recetas con firma","Visitas GPS","Soporte email"]),
    ("Profesional", "Hasta 50 profesionales", C["blue"],
     ["Todo lo de Básico","Emergencias + triage","Telemedicina + app",
      "RRHH + fichajes","Caja + libro diario","Farmacopea + alertas",
      "Chatbot IA","Soporte prioritario"]),
    ("Enterprise", "Ilimitado", C["violet"],
     ["Todo lo de Profesional","Facturación electrónica","API integración",
      "Auditoría legal dedicada","Capacitación presencial","Soporte 24/7"])]):
    x = Inches(0.6+i*4.2)
    RR(s, x, Inches(1.8), Inches(3.9), Inches(5.2), C["card"], color)
    t = TX(s, x+Inches(0.3), Inches(2.0), Inches(3.3), Inches(0.7))
    T(t, nom, 28, True, color, PP_ALIGN.CENTER); T(t, lim, 14, False, C["g"], PP_ALIGN.CENTER)
    R(s, x+Inches(0.5), Inches(3.2), Inches(2.9), Inches(0.02), C["g"])
    t2 = TX(s, x+Inches(0.3), Inches(3.4), Inches(3.3), Inches(3.4))
    for it in items: B(t2, it, 14, RGBColor(200,210,225))
R(s,0,Inches(7.44),W,Inches(0.06),C["gold"])

# ─── 9: WORKFLOW ───────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6]); R(s,0,0,W,H,C["d1"]); R(s,0,0,W,Inches(0.06),C["blue"])
tx = TX(s, Inches(0.8), Inches(0.4), Inches(11), Inches(1))
T(tx, "¿Cómo funciona?", 38, True, C["w"])
T(tx, "De la implementación a la operación diaria", 18, False, C["blue_l"])
R(s, Inches(0.8), Inches(1.5), Inches(3), Inches(0.04), C["blue"])
for i,(paso,tit,desc) in enumerate([
    ("1", "Contacto y demo", "Agendamos una videollamada para conocer sus necesidades y mostrar la plataforma."),
    ("2", "Implementación", "Cargamos pacientes, profesionales y farmacopea. Configuramos roles y permisos."),
    ("3", "Capacitación", "Capacitación guiada por videollamada para todo el equipo."),
    ("4", "Operación diaria", "El equipo usa la plataforma. Soporte continuo por WhatsApp y email.")]):
    y = Inches(1.8+i*1.3)
    RR(s, Inches(0.8), y, Inches(12), Inches(1.1), C["card"], C["line"])
    O(s, Inches(1.2), y+Inches(0.15), Inches(0.8), Inches(0.8), C["teal"])
    t = TX(s, Inches(1.2), y+Inches(0.2), Inches(0.8), Inches(0.7))
    T(t, paso, 24, True, C["w"], PP_ALIGN.CENTER)
    t2 = TX(s, Inches(2.3), y+Inches(0.15), Inches(10), Inches(0.8))
    T(t2, tit, 18, True, C["w"]); T(t2, desc, 14, False, C["g"])
R(s,0,Inches(7.44),W,Inches(0.06),C["blue"])

# ─── 10: CONTACTO ──────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6]); R(s,0,0,W,H,C["d1"]); R(s,0,0,W,Inches(0.06),C["teal"])
O(s,Inches(-1),Inches(2),Inches(6),Inches(6),C["teal"])
O(s,Inches(9.5),Inches(-1),Inches(5),Inches(5),C["blue"])
logo_c(s, Inches(5.9), Inches(0.5), Inches(1.5))
tx = TX(s, Inches(1), Inches(2.2), Inches(11), Inches(1.2))
T(tx, "¿Listo para dar el siguiente paso?", 40, True, C["w"], PP_ALIGN.CENTER)
T(tx, "Agendemos una demo guiada sin compromiso", 22, False, C["teal_l"], PP_ALIGN.CENTER)
R(s, Inches(5.2), Inches(3.6), Inches(3), Inches(0.04), C["teal"])
tx2 = TX(s, Inches(2.5), Inches(4.0), Inches(8.5), Inches(3))
T(tx2, "Enzo N. Girardi", 26, True, C["w"], PP_ALIGN.CENTER)
T(tx2, "Desarrollo técnico y soporte", 15, False, C["g"], PP_ALIGN.CENTER)
T(tx2, "enzogirardi84@gmail.com     |     +54 9 358 430 2024", 16, False, C["g2"], PP_ALIGN.CENTER)
T(tx2, "", 8, False, C["w"])
T(tx2, "Dario Lanfranco", 20, True, C["w"], PP_ALIGN.CENTER)
T(tx2, "Implementación y contratos", 15, False, C["g"], PP_ALIGN.CENTER)
T(tx2, "dariolanfrancoruffener@gmail.com     |     +54 9 358 420 1263", 16, False, C["g2"], PP_ALIGN.CENTER)
R(s,0,Inches(7.44),W,Inches(0.06),C["teal"])

# ─── GUARDAR ───────────────────────────────────────────────────────
prs.save(str(PPTX))
print(f"[OK] PPTX: {PPTX}")
print(f"      {os.path.getsize(str(PPTX))/1024:.1f} KB -- {len(prs.slides)} slides")

# Exportar PDF
TMP = str(Path(tempfile.mktemp(suffix='.pptx')))
try:
    import win32com.client
    prs.save(TMP)  # copia para evitar bloqueo
    ppt = win32com.client.Dispatch("PowerPoint.Application")
    ppt.Visible = True
    pres = ppt.Presentations.Open(TMP, WithWindow=False)
    pres.ExportAsFixedFormat(str(PDF), 2, Intent=1, DocStructureTags=False)
    pres.Close(); ppt.Quit()
    os.remove(TMP)
    print(f"[OK] PDF: {PDF}")
    print(f"      {os.path.getsize(str(PDF))/1024:.1f} KB")
except Exception as e:
    print(f"[WARN] PDF no generado: {e}")
    print("        Abrir PPTX en PowerPoint y guardar como PDF.")
